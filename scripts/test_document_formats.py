from __future__ import annotations

import argparse
import json
import os
import re
import socket
import subprocess
import sys
import time
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Callable

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

from app.extractors import _document_converter, extract_document, split_sections


MARKERS = {
    "txt": "TXT_PIXELRPG_MISTFALL",
    "md": "MD_PIXELRPG_MISTFALL",
    "csv": "CSV_PIXELRPG_MISTFALL",
    "html": "HTML_PIXELRPG_MISTFALL",
    "docx": "DOCX_PIXELRPG_MISTFALL",
    "pptx": "PPTX_PIXELRPG_MISTFALL",
    "xlsx": "XLSX_PIXELRPG_MISTFALL",
    "pdf": "PDF_PIXELRPG_MISTFALL",
    "png": "PNG_PIXELRPG_MISTFALL",
}


@dataclass(slots=True)
class FormatResult:
    extension: str
    bytes: int
    extracted_characters: int
    marker_found: bool
    warnings: list[str]
    seconds: float
    offline_marker_found: bool = False
    offline_warnings: list[str] | None = None


def _minimal_pdf(path: Path, marker: str) -> None:
    stream = f"BT /F1 18 Tf 72 720 Td ({marker}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    data = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, 1):
        offsets.append(len(data))
        data.extend(f"{index} 0 obj\n".encode("ascii") + obj + b"\nendobj\n")
    xref = len(data)
    data.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    data.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        data.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    data.extend(f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode("ascii"))
    path.write_bytes(data)


def _font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for candidate in (Path("C:/Windows/Fonts/arial.ttf"), Path("C:/Windows/Fonts/msjh.ttc")):
        if candidate.is_file():
            return ImageFont.truetype(str(candidate), size)
    return ImageFont.load_default()


def create_fixtures(root: Path) -> dict[Path, str]:
    root.mkdir(parents=True, exist_ok=True)
    fixtures: dict[Path, str] = {}

    for extension, marker, content in (
        ("txt", "TXT_PIXELRPG_MISTFALL", "TXT_PIXELRPG_MISTFALL\n霧落農歌純文字設定。"),
        ("md", "MD_PIXELRPG_MISTFALL", "# MD_PIXELRPG_MISTFALL\n\n霧落農歌 Markdown 設定。"),
        ("csv", "CSV_PIXELRPG_MISTFALL", "id,name\nCSV_PIXELRPG_MISTFALL,霧落作物\n"),
        ("html", "HTML_PIXELRPG_MISTFALL", "<html><body><h1>HTML_PIXELRPG_MISTFALL</h1><p>霧落世界</p></body></html>"),
    ):
        path = root / f"fixture.{extension}"
        path.write_text(content, encoding="utf-8")
        fixtures[path] = marker

    docx_path = root / "fixture.docx"
    doc = Document()
    doc.add_heading("DOCX_PIXELRPG_MISTFALL", level=1)
    doc.add_paragraph("霧落農歌 Word 世界觀文件。")
    doc.save(docx_path)
    fixtures[docx_path] = "DOCX_PIXELRPG_MISTFALL"

    pptx_path = root / "fixture.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PPTX_PIXELRPG_MISTFALL"
    slide.placeholders[1].text = "霧落農歌簡報設定"
    presentation.save(pptx_path)
    fixtures[pptx_path] = "PPTX_PIXELRPG_MISTFALL"

    xlsx_path = root / "fixture.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "PixelRPG"
    sheet.append(["id", "description"])
    sheet.append(["XLSX_PIXELRPG_MISTFALL", "霧落農歌試算表設定"])
    workbook.save(xlsx_path)
    fixtures[xlsx_path] = "XLSX_PIXELRPG_MISTFALL"

    pdf_path = root / "fixture.pdf"
    _minimal_pdf(pdf_path, "PDF_PIXELRPG_MISTFALL")
    fixtures[pdf_path] = "PDF_PIXELRPG_MISTFALL"

    png_path = root / "fixture.png"
    image = Image.new("RGB", (960, 240), "white")
    draw = ImageDraw.Draw(image)
    draw.text((32, 72), "PNG_PIXELRPG_MISTFALL", fill="black", font=_font(46))
    image.save(png_path)
    fixtures[png_path] = "PNG_PIXELRPG_MISTFALL"
    return fixtures


def fixture_map(root: Path) -> dict[Path, str]:
    return {root / f"fixture.{extension}": marker for extension, marker in MARKERS.items()}


def block_external_network() -> Callable[[], None]:
    original_connect = socket.socket.connect

    def guarded_connect(sock: socket.socket, address) -> None:
        host = str(address[0]) if isinstance(address, tuple) and address else ""
        if host not in {"127.0.0.1", "::1", "localhost"}:
            raise RuntimeError(f"Offline document test blocked network connection to {host}")
        return original_connect(sock, address)

    socket.socket.connect = guarded_connect  # type: ignore[method-assign]
    return lambda: setattr(socket.socket, "connect", original_connect)


def extract_all(fixtures: dict[Path, str]) -> list[FormatResult]:
    results: list[FormatResult] = []
    for path, marker in fixtures.items():
        started = time.perf_counter()
        text, warnings = extract_document(path)
        normalized_text = re.sub(r"[^A-Z0-9]+", "", text.upper())
        normalized_marker = re.sub(r"[^A-Z0-9]+", "", marker.upper())
        results.append(FormatResult(
            extension=path.suffix.lower(),
            bytes=path.stat().st_size,
            extracted_characters=len(text),
            marker_found=normalized_marker in normalized_text,
            warnings=warnings,
            seconds=round(time.perf_counter() - started, 3),
        ))
        if text:
            assert split_sections(text), f"{path.suffix} did not produce a chunk"
    return results


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--work-directory", type=Path, required=True)
    parser.add_argument("--report-directory", type=Path, required=True)
    parser.add_argument("--offline-child", action="store_true")
    parser.add_argument("--offline-output", type=Path)
    args = parser.parse_args()
    if args.offline_child:
        if args.offline_output is None:
            raise SystemExit("--offline-output is required for offline child mode")
        restore = block_external_network()
        try:
            offline_results = extract_all(fixture_map(args.work_directory))
        finally:
            restore()
        args.offline_output.write_text(
            json.dumps([asdict(result) for result in offline_results], ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return 0 if all(result.marker_found and not result.warnings for result in offline_results) else 1

    fixtures = create_fixtures(args.work_directory)
    results = extract_all(fixtures)

    offline_output = args.work_directory / "offline-results.json"
    offline_environment = os.environ.copy()
    offline_environment["HF_HUB_OFFLINE"] = "1"
    offline_environment["TRANSFORMERS_OFFLINE"] = "1"
    offline_environment["HF_HUB_DISABLE_TELEMETRY"] = "1"
    offline_process = subprocess.run(
        [
            sys.executable, str(Path(__file__).resolve()),
            "--work-directory", str(args.work_directory),
            "--report-directory", str(args.report_directory),
            "--offline-child", "--offline-output", str(offline_output),
        ],
        env=offline_environment,
        capture_output=True,
        text=True,
        timeout=300,
        check=False,
    )
    offline_values = json.loads(offline_output.read_text(encoding="utf-8")) if offline_output.is_file() else []
    offline_by_extension = {value["extension"]: value for value in offline_values}
    for result in results:
        offline_result = offline_by_extension.get(result.extension, {})
        result.offline_marker_found = bool(offline_result.get("marker_found", False))
        result.offline_warnings = list(offline_result.get("warnings", []))
        if offline_process.returncode != 0 and not result.offline_warnings:
            result.offline_warnings = [f"offline child exit {offline_process.returncode}"]

    passed = all(
        result.marker_found and not result.warnings
        and result.offline_marker_found and not result.offline_warnings
        for result in results
    )
    report = {
        "passed": passed,
        "formats": [asdict(result) for result in results],
        "format_count": len(results),
        "offline_external_connections_allowed": 0,
        "offline_process_exit_code": offline_process.returncode,
    }
    args.report_directory.mkdir(parents=True, exist_ok=True)
    (args.report_directory / "report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    rows = [
        "# PixelRPG 文件格式與離線解析驗收", "",
        f"結果：**{'PASS' if passed else 'FAIL'}**　｜　{sum(r.marker_found for r in results)}/{len(results)} 格式標記擷取　｜　離線外連 0", "",
        "| 格式 | Bytes | 字元 | 首輪 | 離線 | 警告 | 秒 |", "|---|---:|---:|---:|---:|---|---:|",
    ]
    for result in results:
        rows.append(
            f"| {result.extension} | {result.bytes} | {result.extracted_characters} | "
            f"{'通過' if result.marker_found else '失敗'} | {'通過' if result.offline_marker_found else '失敗'} | "
            f"{'；'.join(result.warnings + (result.offline_warnings or []))} | {result.seconds:.3f} |"
        )
    (args.report_directory / "REPORT.md").write_text("\n".join(rows) + "\n", encoding="utf-8")
    print(f"PixelRPG document format gate: {'PASS' if passed else 'FAIL'} ({len(results)} formats, offline rerun)")
    return 0 if passed else 1


if __name__ == "__main__":
    raise SystemExit(main())
